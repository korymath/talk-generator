import pickle
import random
import pathlib
import os.path
import inflect
import argparse
import requests
import safygiphy
from os import listdir
from os.path import isfile, join

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from bs4 import BeautifulSoup

import nltk
from nltk.corpus import wordnet as wn
from py_thesaurus import Thesaurus
from google_images_download import google_images_download

# CONSTANTS
HEIGHT = 9
WIDTH = 16
LEFTMOST = Inches(0)
TOPMOST = Inches(0)
HEIGHT_IN = Inches(HEIGHT)
WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs 
INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
CMS_TO_EMU = 360000


# HELPER FUNCTIONS
def _save_presentation_to_pptx(args, prs):
    """Save the talk."""
    fp = './output/' + args.topic + '.pptx'
    # Create the parent folder if it doesn't exist
    pathlib.Path(os.path.dirname(fp)).mkdir(parents=True, exist_ok=True)
    prs.save(fp)
    print('Saved talk to {}'.format(fp))
    return True


def download_image(fromUrl, toUrl):
    """Download image from url to path."""
    # Create the parent folder if it doesn't exist
    pathlib.Path(os.path.dirname(toUrl)).mkdir(parents=True, exist_ok=True)

    # Download
    f = open(toUrl, 'wb')
    f.write(requests.get(fromUrl).content)
    f.close()


# CONTENT GENERATORS
# These functions generate content, sometimes related to given arguments

def get_definitions(word):
    """Get definitions of a given topic word."""
    print('******************************************')
    # Get definition
    word_senses = wn.synsets(word)
    definitions = {}
    for ss in word_senses:
        definitions[ss.name()] = ss.definition()
    print('{} definitions for "{}"'.format(len(definitions), word))
    return definitions


def get_synonyms(word):
    """Get all synonyms for a given word."""
    print('******************************************')
    word_senses = wn.synsets(word)
    all_synonyms = []
    for ss in word_senses:
        all_synonyms.extend(
            [x.lower().replace('_', ' ') for x in ss.lemma_names()])
    all_synonyms.append(word)
    all_synonyms = list(set(all_synonyms))
    print('{} synonyms for "{}"'.format(len(all_synonyms), word))
    return all_synonyms


def get_relations(word):
    """Get relations to given definitions."""
    rels = {}
    all_rel_forms = []
    all_perts = []
    all_ants = []

    word_senses = wn.synsets(word)
    for ss in word_senses:
        ss_name = ss.name()
        rels[ss_name] = {}
        for lem in ss.lemmas():
            lem_name = lem.name()
            rels[ss_name][lem_name] = {}
            rel_forms = [x.name() for x in lem.derivationally_related_forms()]
            rels[ss_name][lem_name]['related_forms'] = rel_forms
            all_rel_forms.extend(rel_forms)

            perts = [x.name() for x in lem.pertainyms()]
            rels[ss_name][lem_name]['pertainyms'] = perts
            all_perts.extend(perts)

            ants = [x.name() for x in lem.antonyms()]
            rels[ss_name][lem_name]['antonyms'] = ants
            all_ants.extend(ants)

    print('******************************************')
    print('{} derivationally related forms'.format(len(all_rel_forms)))
    print('******************************************')
    print('{} pertainyms'.format(len(all_perts)))
    print('******************************************')
    print('{} antonyms'.format(len(all_ants)))
    return rels


def get_title(synonyms):
    """Returns a template title from a source list."""
    print('******************************************')
    chosen_synonym = random.choice(synonyms)
    chosen_synonym_plural = inflect.engine().plural(chosen_synonym)
    synonym_templates = ['The Unexpected Benefits of {}',
                         'What Your Choice in {} Says About You',
                         'How to Get Rid of {}',
                         'Why {} Will Ruin Your Life',
                         'The Biggest Concerns About {}']
    chosen_template = random.choice(synonym_templates);
    return chosen_template.format(chosen_synonym_plural.title())


def get_images(synonyms, num_images):
    """Get images, first search locally then Google Image Search."""
    all_paths = {}
    if num_images > 0:
        for word in synonyms:
            lp = 'downloads/' + word + '/'
            try:
                local_files = [lp + f for f in listdir(lp) if isfile(join(lp, 
                    f))]
                all_paths[word] = local_files
            except FileNotFoundError as e:
                all_paths[word] = []
                pass

            if len(all_paths[word]) > 0:
                print('{} local images on {} found'.format(len(all_paths[word]),
                                                           word))
            # If no local images, search on Google Image Search
            if len(all_paths[word]) == 0:
                # Get related images at 16x9 aspect ratio
                response = google_images_download.googleimagesdownload()
                arguments = {
                    'keywords': word,
                    'limit': num_images,
                    'print_urls': True,
                    'exact_size': '1600,900',
                }
                # passing the arguments to the function
                paths = response.download(arguments)
                # printing absolute paths of the downloaded images
                print('paths of images', paths)
                # Add to main dictionary
                all_paths[word] = paths[word]
    return all_paths


def get_related_giphy(seed_word):
    giphy = safygiphy.Giphy()
    result = giphy.random(tag=seed_word)
    return result.get('data').get('images').get('original').get('url')


def wikihow_action_to_action(wikihow_title):
    index_of_to = wikihow_title.find('to')
    return wikihow_title[index_of_to + 3:]


def get_related_wikihow_actions(seed_word):
    page = requests.get(('https://en.wikihow.com/'
        'wikiHowTo?search=') + seed_word.replace(' ', '+'))
    # Try again but with plural if nothing is found
    if not page:
        page = requests.get(('https://en.wikihow.com/wikiHowTo?search=' 
            + inflect.engine().plural(seed_word).replace(" ", "+")))


    soup = BeautifulSoup(page.content, 'html.parser')
    actions_elements = soup.find_all('a', class_='result_link')
    actions = \
        list(
            map(wikihow_action_to_action,
                map(lambda x: x.get_text(), actions_elements)))

    #print(actions)

    return actions


# FORMAT GENERATORS
# These are functions that get some inputs (texts, images...) 
# and create layouted slides with these inputs

def create_title_slide(args, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_object = slide.shapes.title
    title_object.text = args.title
    title_object.width = WIDTH_IN
    title_object.height = HEIGHT_IN
    title_object.left = LEFTMOST
    title_object.right = TOPMOST
    return slide


def create_text_slide(prs, text):
    # Get a default blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_object = slide.shapes.title
    title_object.text = text
    title_object.width = WIDTH_IN
    title_object.height = HEIGHT_IN
    title_object.left = LEFTMOST
    title_object.right = TOPMOST
    return slide


def create_image_slide(prs, image_url):
    # Get a default blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add image url as picture
    if image_url:
        pic = slide.shapes.add_picture(image_url,
                                       LEFTMOST, TOPMOST, height=HEIGHT_IN)
        return slide
    return False


# FULL SLIDES GENERATORS:
# These are functions that create slides with certain (generated) content

def create_google_image_slide(args, prs, word):
    # Get all image paths
    img_paths = args.all_paths.get(word)
    if img_paths:
        # Pick one of the images
        img_path = random.choice(img_paths)

        # Create slide with image
        slide = create_image_slide(prs, img_path)

        # Add title to the slide
        if slide:
            shapes = slide.shapes
            shapes.title.text = word
            return slide
    return False


def create_inspirobot_slide(prs):
    # Generate a random url to access inspirobot
    dd = str(random.randint(1, 73)).zfill(2)
    nnnn = random.randint(0, 9998)
    inspirobot_url = ('http://generated.inspirobot.me/'
                      '0{}/aXm{}xjU.jpg').format(dd, nnnn)

    # Download the image
    image_url = 'downloads/inspirobot/{}-{}.jpg'.format(dd, nnnn)
    download_image(inspirobot_url, image_url)

    print("Downloaded inspirobot image: {}".format(image_url))

    # Turn into image slide
    return create_image_slide(prs, image_url)


def create_giphy_slide(prs, word):
    # Download the image
    giphy_url = get_related_giphy(word)
    gif_name = os.path.basename(os.path.dirname(giphy_url))
    image_url = 'downloads/' + word + '/gifs/' + gif_name + ".gif"
    download_image(giphy_url, image_url)

    # Turn into image slide
    return create_image_slide(prs, image_url)


def create_wikihow_action_recommendation_slide(prs, wikihow_seed):
    related_actions = get_related_wikihow_actions(wikihow_seed)
    if related_actions:
        action = random.choice(related_actions)
        life_lesson_templates = ['Life hack: Always {}!',
                                 'I Will Teach You How To {}!',
                                 'Life Advice: {}!',
                                 'Life Advice: Never {}!',
                                 'WARNING: Never {}!',
                                 'Friendly Reminder to {}',
                                 'When in Doubt: {}']

        life_lesson = random.choice(
            life_lesson_templates).format(action.title())

        # Turn into image slide
        return create_text_slide(prs, life_lesson)
    return False


# COMPILATION
# Compiling the slides to a powerpoint


def compile_talk_to_pptx(args):
    """Compile the talk with the given source material."""
    prs = Presentation()
    # Set the height and width
    prs.slide_height = HEIGHT * INCHES_TO_EMU
    prs.slide_width = WIDTH * INCHES_TO_EMU

    # Build an ordered list of slides for access
    slides = []

    # Add title slide
    title_slide = create_title_slide(args, prs)
    slides.append(title_slide)
    slide_idx_iter = 1

    # For each synonym 
    for word, path_list in args.all_paths.items():
        # print('Word: {}'.format(word))
        # For each image collected add a new slide
        # for i in range(len(path_list)):
        print('***********************************')
        print('Adding slide {} about {}'.format(slide_idx_iter, word))
        slide = create_google_image_slide(args, prs, word)
        if slide:
            slides.append(slide)
            slide_idx_iter += 1

    # Add some Inspirobot quotes
    print('***********************************')
    print('Adding slide: {}, Inspirobot'.format(slide_idx_iter))
    slide = create_inspirobot_slide(prs)
    if slide:
        slides.append(slide)
        slide_idx_iter += 1

    # Add a Gif slide
    print('***********************************')
    giphy_seed = random.choice(args.synonyms)
    print('Adding slide: {}, Giphy about {}'.format(slide_idx_iter, giphy_seed))
    slide = create_giphy_slide(prs, giphy_seed)
    if slide:
        slides.append(slide)
        slide_idx_iter += 1

    # Add a life lesson
    print('***********************************')
    wikihow_seed = random.choice(args.synonyms)
    print('Adding Wikihow Lifelesson slide: {} about {}'.format(slide_idx_iter, 
        giphy_seed))
    slide = create_wikihow_action_recommendation_slide(prs, wikihow_seed)
    if slide:
        slides.append(slide)
        slide_idx_iter += 1

    _save_presentation_to_pptx(args, prs)
    print('Successfully built talk.')


def compile_talk_to_raw_data(args):
    """Save the raw data that has been harvested for future use."""
    return True


# MAIN

def main(args):
    """Make a talk with the given topic."""
    # Print status details
    print('******************************************')
    print("Making {} slide talk on: {}".format(args.num_slides, args.topic))

    # Parse topic string to parts-of-speech
    text = nltk.word_tokenize(args.topic)
    print('tokenized text: ', text)
    print('pos tag text: ', nltk.pos_tag(text))

    # Parse the actual topic subject from the parts-of-speech
    topic_string = args.topic

    # Get definitions
    args.definitions = get_definitions(topic_string)
    # Get relations
    args.relations = get_relations(topic_string)
    # Get synonyms
    args.synonyms = get_synonyms(topic_string)
    # Get related actions
    args.actions = get_related_wikihow_actions(topic_string)
    # Get a title
    args.title = get_title(args.synonyms)
    # For each synonym download num_images
    args.all_paths = get_images(args.synonyms, args.num_images)

    # Compile and save the presentation to data
    compile_talk_to_raw_data(args)

    # Compile and save the presentation to PPTX
    compile_talk_to_pptx(args)


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('--topic', help="Topic of presentation.",
                        default='bagels', type=str)
    parser.add_argument('--num_images', help="Number of images per synonym.",
                        default=1, type=int)
    parser.add_argument('--num_slides', help="Number of slides to create.",
                        default=3, type=int)
    args = parser.parse_args()
    main(args)
