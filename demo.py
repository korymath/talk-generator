import argparse


from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from nltk.corpus import wordnet as wn
from py_thesaurus import Thesaurus
from PyDictionary import PyDictionary
from google_images_download import google_images_download

# CONSTANTS
HEIGHT = 9
WIDTH = 16
LEFTMOST = Inches(0)
TOPMOST = Inches(0)
HEIGHT_IN = Inches(HEIGHT)
WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs 
INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
CMS_TO_EMU = 360000


def get_definitions(args):
  print('******************************************')
  print('Topic: {}'.format(args.topic))
  
  # Get definition
  dictionary = PyDictionary()
  definitions = dictionary.meaning(args.topic)
  
  if definitions:
    print('******************************************')
    print('{} word type(s)'.format(len(definitions)))
    for word_type,word_type_defs in definitions.items():
      print('******************************************')
      print('As a {}, {} can mean {} things: '.format(word_type.lower(), 
        args.topic, len(word_type_defs)))
      
      for word_type_def in word_type_defs:
        print('\t {}'.format(word_type_def))
    return definitions
  else:
    print('No definition found.')
    return None

def get_synonyms(args):
  print('******************************************')
  # Get N synonyms
  word_senses = wn.synsets(args.topic)

  all_synonyms = []
  for ss in word_senses:
    # print(ss.name(), ss.lemma_names(), ss.definition())
    all_synonyms.extend([x.lower().replace('_', ' ') for x in ss.lemma_names()])

  all_synonyms = list(set(all_synonyms))

  print('{} synonyms: '.format(len(all_synonyms)))
  for synonym in all_synonyms:
    if synonym is not args.topic.lower():
      print('\t {}'.format(synonym))

  return all_synonyms

def get_images(synonyms, num_images):
  all_paths = {}
  for synonym in synonyms:
    # Get related images at 16x9 aspect ratio

    # TODO: add image filter for weird and NSFW stuff
    response = google_images_download.googleimagesdownload()
    arguments = {
      'keywords':synonym,
      'limit':num_images,
      'print_urls':True,
      'exact_size':'1600,900',
      # 'size':'large',
      # 'usage_rights':'labeled-for-noncommercial-reuse-with-modification'
    }
    # passing the arguments to the function
    paths = response.download(arguments)
    # printing absolute paths of the downloaded images
    print(paths)
    # Add to main dictionary
    all_paths[synonym] = paths[synonym]
  return all_paths

def compile_presentation(args, all_paths, definitions, synonyms):
  # Make a presentation
  prs = Presentation()
  
  # Set the height and width
  prs.slide_height = HEIGHT * INCHES_TO_EMU
  prs.slide_width = WIDTH * INCHES_TO_EMU
  
  # Get a default blank slide layout
  slide_layout = prs.slide_layouts[5]

  # Build an ordered list of slides for access
  slides = []

  slide_idx_iter = 0
  for synonym,paths in all_paths.items():
    print('***********************************')
    print('Adding slide: {}'.format(slide_idx_iter))
    slides.append(prs.slides.add_slide(slide_layout))
    
    # Get an image
    img_path = paths[0]

    # Add the image to the slide.
    pic = slides[slide_idx_iter].shapes.add_picture(img_path, LEFTMOST, TOPMOST, 
      width=WIDTH_IN, height=HEIGHT_IN)

    # Add title to the slide
    shapes = slides[slide_idx_iter].shapes
    shapes.title.text = synonym

    # TODO: Add the text to the slide.

    slide_idx_iter += 1

  return prs,slides

def save_talk(args, prs):
  # Save the presentation
  fp = 'output/' + args.topic + '-' + args.output
  # TODO: make directory if it doesn't exist
  # FileNotFoundError: [Errno 2] No such file or directory: 'output/latitude-test.pptx'
  prs.save(fp)
  print('Saved talk to {}'.format(fp))
  return True

def main(args):
  """Make a presentation with the given topic.""" 

  # Get definitions
  definitions = get_definitions(args)
  synonyms = get_synonyms(args)

  # For each synonym get N image paths
  all_paths = get_images(synonyms, args.num_images)

  # Compile the presentation
  prs, slides = compile_presentation(args, all_paths=all_paths, 
    definitions=definitions, synonyms=synonyms)

  if save_talk(args, prs):
    print('Successfully built talk.')

if __name__ == '__main__':
  parser = argparse.ArgumentParser(description='Process some integers.')
  parser.add_argument('--output', help="Output filename.", 
    default='test.pptx', type=str)
  parser.add_argument('--topic', help="Topic of presentation.", 
    default='bagels', type=str)
  parser.add_argument('--num_images', help="Number of images for each synonym.", 
    default=1, type=int)
  parser.add_argument('--num_slides', help="Number of slides to create.", 
    default=6, type=int)
  args = parser.parse_args()
  main(args)



  